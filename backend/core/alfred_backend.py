"""
Alfred Backend API - FastAPI Server
"""

import sys
import io
import os
from pathlib import Path

# Agregar directorios al path de Python para imports
backend_root = Path(__file__).parent.parent
sys.path.insert(0, str(backend_root))
sys.path.insert(0, str(backend_root / "core"))
sys.path.insert(0, str(backend_root / "gpu"))
sys.path.insert(0, str(backend_root / "utils"))

from utils.logger import get_logger

# ============================================
# AUTO-REPARACION DE DEPENDENCIAS
# ============================================
try:
    from utils.auto_repair import run_auto_repair, get_repair_status
    
    logger_temp = get_logger("auto_repair")
    logger_temp.info("[STARTUP] Verificando integridad del sistema...")
    
    # Verificar estado antes de continuar
    repair_status = get_repair_status()
    
    if repair_status['overall'] == 'needs_repair':
        logger_temp.warning("[STARTUP] Detectados problemas - Aplicando correcciones automaticas...")
        
        # Ejecutar reparacion
        success = run_auto_repair()
        
        if not success:
            logger_temp.warning("[STARTUP] Correcciones aplicadas - Reiniciando automaticamente...")
            # NO imprimir mensaje visible al usuario en produccion
            # Electron detectara el exit code 3 y reiniciara automaticamente
            sys.exit(3)  # Codigo 3 = auto-reparacion completada, reinicio automatico
    else:
        logger_temp.info("[STARTUP] Sistema OK - Continuando inicio...")
        
except Exception as e:
    # Si falla la auto-reparacion, continuar con advertencia
    print(f"[STARTUP WARNING] No se pudo ejecutar auto-reparacion: {e}")

# ============================================
# CONTINUACION DEL INICIO NORMAL
# ============================================

from db_manager import init_db
import db_manager

# Configurar encoding UTF-8 para evitar errores con caracteres especiales en Windows
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from fastapi import FastAPI, HTTPException, BackgroundTasks, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field, field_validator
from typing import Optional, List, Dict, Any
from datetime import datetime
from contextlib import asynccontextmanager
import asyncio
import json

# Importar módulos locales
import functionsToHistory
from alfred_core import AlfredCore
from conversation_manager import get_conversation_manager
from utils.security import encrypt_data, decrypt_data, encrypt_for_transport, is_encryption_enabled
from functionsToHistory import encrypt_personal_data, decrypt_personal_data

# Importar routers de endpoints
from endpoints.user.user import router as user_router
from endpoints.conversations.conversations import router as conversations_router
from endpoints.documents.documents import router as documents_router
from endpoints.security.security import router as security_router
from endpoints.shared_state import set_alfred_core_instance
# Importar modelos de usuario si se necesitan en otros endpoints
# from endpoints.user.user import UserSettingRequest, UserSettingResponse, ProfilePictureRequest, ProfilePictureHistoryResponse

# --- Utilidades para procesamiento de archivos ---

def extract_text_from_pdf(content: str, file_name: str) -> str:
    """
    Extraer texto de un PDF a partir de su contenido en base64 o texto
    
    Args:
        content: Contenido del PDF (puede ser base64 con o sin prefijo data:)
        file_name: Nombre del archivo para logging
        
    Returns:
        Texto extraído del PDF
    """
    try:
        import pypdf
        import base64
        from io import BytesIO
        import re
        
        print(f"🔍 Procesando PDF: {file_name}")
        print(f"📏 Longitud del contenido recibido: {len(content)} caracteres")
        
        # Limpiar prefijo data: si existe (ej: data:application/pdf;base64,)
        if content.startswith('data:'):
            print(f"🧹 Limpiando prefijo data: URL...")
            # Buscar la coma que separa el header del contenido base64
            match = re.search(r'base64,(.+)', content)
            if match:
                content = match.group(1)
                print(f"✅ Prefijo removido, nueva longitud: {len(content)} caracteres")
            else:
                print(f"⚠️ No se encontro separador base64, usando contenido completo")
        
        # Decodificar base64
        try:
            print(f"🔓 Decodificando base64...")
            pdf_bytes = base64.b64decode(content)
            print(f"✅ Decodificado exitoso: {len(pdf_bytes)} bytes")
        except Exception as decode_error:
            print(f"❌ Error al decodificar base64: {str(decode_error)}")
            raise ValueError(f"Error al decodificar base64: {str(decode_error)}")
        
        # Crear objeto PDF
        print(f"📄 Creando lector PDF...")
        pdf_file = BytesIO(pdf_bytes)
        pdf_reader = pypdf.PdfReader(pdf_file)
        
        print(f"📚 PDF cargado: {len(pdf_reader.pages)} paginas")
        
        # Extraer texto de todas las páginas
        text_parts = []
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Pagina {page_num + 1} ---\n{text}")
                    print(f"  ✅ Pagina {page_num + 1}: {len(text)} caracteres")
                else:
                    print(f"  ⚠️ Pagina {page_num + 1}: vacia o sin texto extraible")
            except Exception as page_error:
                print(f"  ❌ Error en pagina {page_num + 1}: {str(page_error)}")
                text_parts.append(f"--- Pagina {page_num + 1} ---\n[Error al extraer texto de esta pagina]")
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ PDF procesado exitosamente: {file_name}")
        print(f"📊 Total: {len(pdf_reader.pages)} paginas, {len(extracted_text)} caracteres extraidos")
        
        return extracted_text if extracted_text else "[PDF procesado pero no se pudo extraer texto]"
        
    except Exception as e:
        import traceback
        error_detail = traceback.format_exc()
        print(f"❌ Error al procesar PDF {file_name}:")
        print(error_detail)
        return f"[Error al procesar PDF: {str(e)}. Por favor verifica que el archivo no este corrupto.]"


def extract_text_from_docx(content: str, file_name: str) -> str:
    """
    Extraer texto de un documento Word (.docx) a partir de su contenido en base64
    
    Args:
        content: Contenido del archivo en base64 (con o sin prefijo data:)
        file_name: Nombre del archivo para logging
        
    Returns:
        Texto extraído del documento
    """
    try:
        from docx import Document
        import base64
        from io import BytesIO
        import re
        
        print(f"🔍 Procesando Word: {file_name}")
        
        # Limpiar prefijo data: si existe
        if content.startswith('data:'):
            match = re.search(r'base64,(.+)', content)
            if match:
                content = match.group(1)
        
        # Decodificar contenido base64
        docx_bytes = base64.b64decode(content)
        docx_file = BytesIO(docx_bytes)
        
        # Cargar documento
        doc = Document(docx_file)
        
        # Extraer texto de parrafos
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extraer texto de tablas
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n--- Tabla ---\n" + "\n".join(table_text))
        
        extracted_text = "\n\n".join(text_parts)
        print(f"✅ Word procesado: {file_name} - {len(doc.paragraphs)} parrafos, {len(extracted_text)} caracteres")
        
        return extracted_text if extracted_text else "[Documento Word procesado pero no se pudo extraer texto]"
        
    except Exception as e:
        import traceback
        print(f"❌ Error al procesar Word {file_name}:")
        print(traceback.format_exc())
        return f"[Error al procesar documento Word: {str(e)}]"


def extract_text_from_xlsx(content: str, file_name: str) -> str:
    """
    Extraer texto de una hoja de calculo Excel (.xlsx) a partir de su contenido en base64
    
    Args:
        content: Contenido del archivo en base64 (con o sin prefijo data:)
        file_name: Nombre del archivo para logging
        
    Returns:
        Texto extraído de todas las hojas
    """
    try:
        from openpyxl import load_workbook
        import base64
        from io import BytesIO
        import re
        
        print(f"🔍 Procesando Excel: {file_name}")
        
        # Limpiar prefijo data: si existe
        if content.startswith('data:'):
            match = re.search(r'base64,(.+)', content)
            if match:
                content = match.group(1)
        
        # Decodificar contenido base64
        xlsx_bytes = base64.b64decode(content)
        xlsx_file = BytesIO(xlsx_bytes)
        
        # Cargar workbook
        wb = load_workbook(xlsx_file, data_only=True)
        
        # Extraer texto de cada hoja
        sheet_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Obtener valores de celdas
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows_text.append(row_text)
            
            if rows_text:
                sheet_parts.append(f"--- Hoja: {sheet_name} ---\n" + "\n".join(rows_text))
        
        extracted_text = "\n\n".join(sheet_parts)
        print(f"✅ Excel procesado: {file_name} - {len(wb.sheetnames)} hojas, {len(extracted_text)} caracteres")
        
        return extracted_text if extracted_text else "[Hoja de calculo procesada pero no se pudo extraer texto]"
        
    except Exception as e:
        import traceback
        print(f"❌ Error al procesar Excel {file_name}:")
        print(traceback.format_exc())
        return f"[Error al procesar hoja de calculo: {str(e)}]"


def extract_text_from_pptx(content: str, file_name: str) -> str:
    """
    Extraer texto de una presentacion PowerPoint (.pptx) a partir de su contenido en base64
    
    Args:
        content: Contenido del archivo en base64 (con o sin prefijo data:)
        file_name: Nombre del archivo para logging
        
    Returns:
        Texto extraído de todas las diapositivas
    """
    try:
        from pptx import Presentation
        import base64
        from io import BytesIO
        import re
        
        print(f"🔍 Procesando PowerPoint: {file_name}")
        
        # Limpiar prefijo data: si existe
        if content.startswith('data:'):
            match = re.search(r'base64,(.+)', content)
            if match:
                content = match.group(1)
        
        # Decodificar contenido base64
        pptx_bytes = base64.b64decode(content)
        pptx_file = BytesIO(pptx_bytes)
        
        # Cargar presentacion
        prs = Presentation(pptx_file)
        
        # Extraer texto de cada diapositiva
        slide_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extraer texto de todas las formas
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                slide_parts.append(f"--- Diapositiva {slide_num} ---\n" + "\n".join(slide_text))
        
        extracted_text = "\n\n".join(slide_parts)
        print(f"✅ PowerPoint procesado: {file_name} - {len(prs.slides)} diapositivas, {len(extracted_text)} caracteres")
        
        return extracted_text if extracted_text else "[Presentacion procesada pero no se pudo extraer texto]"
        
    except Exception as e:
        import traceback
        print(f"❌ Error al procesar PowerPoint {file_name}:")
        print(traceback.format_exc())
        return f"[Error al procesar presentacion: {str(e)}]"


# --- Modelos de datos para la API ---

class QueryRequest(BaseModel):
    """Solicitud de consulta al asistente con validacion mejorada"""
    question: str = Field(
        ..., 
        description="Pregunta del usuario", 
        min_length=1,
        max_length=2000
    )
    use_history: bool = Field(True, description="Buscar primero en el historial")
    save_response: bool = Field(False, description="Guardar respuesta automaticamente")
    search_documents: bool = Field(True, description="Buscar en documentos o solo usar el prompt")
    search_kwargs: Optional[Dict[str, Any]] = Field(
        None, 
        description="Parametros adicionales de busqueda (k, fetch_k, search_type)"
    )
    
    @field_validator('question')
    @classmethod
    def validate_question(cls, v: str) -> str:
        """Validar que la pregunta no este vacia despues de strip"""
        if not v or not v.strip():
            raise ValueError('La pregunta no puede estar vacia')
        return v.strip()
    
    @field_validator('search_kwargs')
    @classmethod
    def validate_search_kwargs(cls, v: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        """Validar que search_kwargs solo contenga claves permitidas"""
        if v is not None:
            allowed_keys = {'k', 'fetch_k', 'search_type', 'score_threshold'}
            invalid_keys = set(v.keys()) - allowed_keys
            if invalid_keys:
                raise ValueError(
                    f'Claves no permitidas en search_kwargs: {invalid_keys}. '
                    f'Claves permitidas: {allowed_keys}'
                )
        return v

class QueryResponse(BaseModel):
    """Respuesta del asistente"""
    answer: str = Field(..., description="Respuesta generada por Alfred")
    personal_data: Optional[Dict[str, str]] = Field(None, description="Datos personales extraídos")
    sources: List[str] = Field(default_factory=list, description="Fuentes de los documentos")
    from_history: bool = Field(False, description="Si la respuesta proviene del historial")
    history_score: Optional[float] = Field(None, description="Score de similitud con historial")
    timestamp: str = Field(default_factory=lambda: datetime.now().isoformat())
    context_count: int = Field(0, description="Número de fragmentos recuperados")
    from_cache: Optional[bool] = Field(None, description="Si la respuesta proviene del cache en memoria")
    cache_age_seconds: Optional[float] = Field(None, description="Edad del cache en segundos")
    
    class Config:
        extra = "allow"  # Permitir campos adicionales para compatibilidad futura

class HistorySearchRequest(BaseModel):
    """Solicitud de búsqueda en historial"""
    search_term: str = Field(..., description="Término a buscar", min_length=1)
    threshold: float = Field(0.2, description="Umbral de similitud", ge=0.0, le=1.0)
    top_k: int = Field(10, description="Número máximo de resultados", ge=1, le=50)

class HistoryEntry(BaseModel):
    """Entrada del historial"""
    timestamp: str
    question: str
    answer: str
    personal_data: Optional[Dict[str, str]] = None
    sources: List[str] = Field(default_factory=list)
    similarity_score: Optional[float] = None

class OptimizationStats(BaseModel):
    """Estadisticas de optimizaciones RAG"""
    embedding_model: str = Field(..., description="Modelo de embeddings en uso")
    embedding_dimension: int = Field(..., description="Dimension de los vectores")
    vram_available: float = Field(..., description="VRAM disponible en GB")
    cache_enabled: bool = Field(..., description="Si el cache LRU esta activo")
    cache_hits: int = Field(0, description="Numero de hits del cache")
    cache_misses: int = Field(0, description="Numero de misses del cache")
    cache_hit_rate: float = Field(0.0, description="Tasa de acierto del cache")
    cache_size: int = Field(0, description="Entradas actuales en cache")
    total_documents_indexed: int = Field(0, description="Documentos en vector store")
    chunking_strategies: Dict[str, int] = Field(default_factory=dict, description="Estrategias de chunking aplicadas")
    optimized_storage: bool = Field(False, description="Si usa DuckDB+Parquet")
    storage_path: str = Field(..., description="Ruta del almacenamiento ChromaDB")

class ModelInfo(BaseModel):
    """Información del modelo actual"""
    model_name: str
    available_models: List[str] = Field(default_factory=list)

class ChangeModelRequest(BaseModel):
    """Solicitud para cambiar el modelo"""
    model_name: str = Field(..., description="Nombre del nuevo modelo a utilizar")

class GPUStatus(BaseModel):
    """Estado de la GPU"""
    gpu_available: bool
    device_type: str
    device: str
    gpu_info: Dict[str, Any] = Field(default_factory=dict)
    memory_usage: Optional[Dict[str, float]] = None

class HealthResponse(BaseModel):
    """Estado de salud del servicio con detalles de componentes"""
    status: str = Field(..., description="Estado general: healthy, degraded, unhealthy")
    timestamp: str = Field(..., description="Timestamp del health check")
    components: Dict[str, str] = Field(default_factory=dict, description="Estado de cada componente")
    alfred_core_initialized: bool = Field(..., description="Si Alfred Core esta inicializado")
    vectorstore_loaded: bool = Field(..., description="Si la base vectorial esta cargada")
    gpu_status: Optional[GPUStatus] = Field(None, description="Estado detallado de GPU")
    is_fully_initialized: bool = Field(..., description="Si la inicializacion asincrona ha completado completamente (lifespan ready)")

class OllamaKeepAliveRequest(BaseModel):
    """Solicitud para actualizar el keep_alive de Ollama"""
    seconds: int = Field(..., description="Tiempo en segundos (0-3600)", ge=0, le=3600)

class OllamaKeepAliveResponse(BaseModel):
    """Respuesta con el keep_alive actual"""
    keep_alive_seconds: int = Field(..., description="Tiempo en segundos")
    description: str = Field(..., description="Descripcion del comportamiento")

# --- Modelos de Modo de Aplicacion ---

class ModeRequest(BaseModel):
    """Solicitud para cambiar el modo de la aplicacion"""
    mode: str = Field(..., description="Modo: work, focus, personal, creative")
    
    @field_validator('mode')
    @classmethod
    def validate_mode(cls, v: str) -> str:
        """Validar que el modo sea valido"""
        allowed_modes = ['work', 'focus', 'personal', 'creative']
        if v not in allowed_modes:
            raise ValueError(f'Modo invalido. Modos permitidos: {allowed_modes}')
        return v

class ModeResponse(BaseModel):
    """Respuesta con el modo actual"""
    mode: str = Field(..., description="Modo actual")
    updated_at: Optional[str] = Field(None, description="Fecha de ultima actualizacion")

# --- Modelos de Tema ---

class ThemeRequest(BaseModel):
    """Solicitud para cambiar el tema de la aplicacion"""
    theme: str = Field(..., description="Tema: light, dark")
    
    @field_validator('theme')
    @classmethod
    def validate_theme(cls, v: str) -> str:
        """Validar que el tema sea valido"""
        allowed_themes = ['light', 'dark']
        if v not in allowed_themes:
            raise ValueError(f'Tema invalido. Temas permitidos: {allowed_themes}')
        return v

class ThemeResponse(BaseModel):
    """Respuesta con el tema actual"""
    theme: str = Field(..., description="Tema actual")
    updated_at: Optional[str] = Field(None, description="Fecha de ultima actualizacion")

# --- Inicialización del núcleo de Alfred ---
alfred_core: Optional[AlfredCore] = None
_initialization_complete: bool = False  # Flag para indicar cuando initialize_async() ha terminado

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Manejo del ciclo de vida de la aplicación"""
    global alfred_core
    global _initialization_complete
    
    # IMPORTANTE: Usar sys.stdout.flush() para forzar salida inmediata en Windows
    import sys
    
    print("\n" + "="*60, flush=True)
    print("Iniciando Alfred Backend API...", flush=True)
    print("="*60, flush=True)
    sys.stdout.flush()
    
    _initialization_complete = False  # Marcar como no listo
    
    try:
        # Inicializar el núcleo de Alfred con version refactorizada (async + optimizaciones)
        print("Creando instancia de AlfredCore...", flush=True)
        alfred_core = AlfredCore()
        
        print("Inicializando componentes async (lazy loading)...", flush=True)
        await alfred_core.initialize_async()
        
        # Establecer instancia global para endpoints modulares
        set_alfred_core_instance(alfred_core)
        
        print("\n" + "="*60, flush=True)
        print("Alfred Core Refactored inicializado correctamente", flush=True)
        print(f"  - Embedding Model: {alfred_core.embedding_model}", flush=True)
        print(f"  - Vector Store: {alfred_core.chroma_db_path}", flush=True)
        print(f"  - Optimizations: Incremental indexing + LRU cache + Adaptive chunking", flush=True)
        print("="*60 + "\n", flush=True)
        sys.stdout.flush()
        
        # MARCAR COMO COMPLETAMENTE LISTO SOLO DESPUES DE initialize_async()
        _initialization_complete = True
        print("INITIALIZATION COMPLETE: Backend completamente listo para procesar consultas", flush=True)
        sys.stdout.flush()
        
        # Establecer instancia compartida de alfred_core para endpoints modulares
        set_alfred_core_instance(alfred_core)
        
        yield
    except Exception as e:
        print(f"\nError al inicializar Alfred Core: {e}", flush=True)
        import traceback
        traceback.print_exc()
        sys.stdout.flush()
        _initialization_complete = False  # Marcar como fallido
        raise
    finally:
        print("\n" + "="*60, flush=True)
        print("Cerrando Alfred Backend API...", flush=True)
        print("="*60 + "\n", flush=True)
        sys.stdout.flush()
        # Aquí podrías agregar limpieza si es necesaria


backend_logger = get_logger("server")
rag_logger = get_logger("rag")
security_logger = get_logger("security")

# Log de inicio
backend_logger.info(f"Iniciando backend Alfred en {os.getenv('ALFRED_IP', 'Not found')}:{os.getenv('ALFRED_PORT', 'Not found')}")
rag_logger.info(f"Iniciando RAG en {os.getenv('ALFRED_RAG_IP', 'Not found')}:{os.getenv('ALFRED_RAG_PORT', 'Not found')}")
security_logger.info(f"Iniciando seguridad en {os.getenv('ALFRED_SECURITY_IP', 'Not found')}:{os.getenv('ALFRED_SECURITY_PORT', 'Not found')}")

# Inicio de la base de datos
init_db()

# --- Funciones auxiliares de cifrado ---

def ensure_personal_data_decrypted(data: Optional[Dict[str, str]]) -> Optional[Dict[str, str]]:
    """
    Asegura que los datos personales esten descifrados antes de enviarlos al cliente
    
    Args:
        data: Diccionario con datos personales (potencialmente cifrados)
    
    Returns:
        Diccionario con datos descifrados o None
    """
    if not data:
        return None
    
    try:
        # Intentar descifrar cada campo
        decrypted = {}
        for key, value in data.items():
            if value and isinstance(value, str):
                try:
                    # Intentar descifrar
                    decrypted[key] = decrypt_data(value)
                    security_logger.debug(f"Campo {key} descifrado exitosamente")
                except Exception:
                    # Si falla, asumir que ya esta descifrado
                    decrypted[key] = value
            else:
                decrypted[key] = value
        return decrypted
    except Exception as e:
        security_logger.error(f"Error al descifrar datos personales: {e}")
        return data  # Devolver datos originales en caso de error

def log_personal_data_access(operation: str, data_keys: List[str], user_context: str = "API"):
    """
    Registra el acceso a datos personales para auditoria
    
    Args:
        operation: Tipo de operacion (read, write, delete)
        data_keys: Claves de los datos accedidos (rfc, curp, etc.)
        user_context: Contexto del usuario/cliente
    """
    security_logger.info(
        f"Acceso a datos personales: {operation} | "
        f"Campos: {', '.join(data_keys)} | "
        f"Contexto: {user_context} | "
        f"Timestamp: {datetime.now().isoformat()}"
    )

# --- Crear aplicación FastAPI ---
app = FastAPI(
    title="Alfred Backend API",
    description="Asistente personal inteligente con acceso a documentos locales",
    version="1.0.0",
    lifespan=lifespan,
    docs_url="/docs",  # Swagger UI
    redoc_url="/redoc"  # ReDoc
)

# --- Configurar CORS para permitir peticiones desde C# ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # En producción, especifica los orígenes permitidos
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Incluir routers de módulos ---

## Router de los endpoints de usuario
app.include_router(user_router)

## Router de los endpoints de conversaciones
app.include_router(conversations_router)

## Router de los endpoints de documentos
app.include_router(documents_router)

## Router de los endpoints de seguridad y cifrado
app.include_router(security_router)

# --- Endpoints ---

@app.get("/", tags=["General"])
async def root():
    """Endpoint raíz - Información del servicio"""
    return {
        "service": "Alfred Backend API",
        "version": "1.0.0",
        "status": "running",
        "docs": "/docs",
        "health": "/health"
    }

@app.get("/health", response_model=HealthResponse, tags=["General"])
async def health_check():
    """
    Verificar el estado de salud del servicio con detalles de componentes
    
    Retorna el estado de:
    - Alfred Core (inicializacion)
    - ChromaDB (base vectorial con conteo de documentos)
    - Ollama (LLM disponible)
    - GPU (si esta disponible y configurada)
    """
    health_status = "healthy"
    components = {}
    gpu_status_data = None
    
    # Check Alfred Core
    if alfred_core and alfred_core.is_initialized():
        components["alfred_core"] = "healthy"
    else:
        health_status = "unhealthy"
        components["alfred_core"] = "not_initialized"
    
    # Check ChromaDB
    try:
        if alfred_core and alfred_core.vector_manager and alfred_core.vector_manager.vectorstore:
            count = alfred_core.vector_manager.vectorstore._collection.count()
            components["chroma_db"] = f"healthy ({count} docs)"
        else:
            health_status = "degraded"
            components["chroma_db"] = "not_available"
    except Exception as e:
        health_status = "degraded"
        components["chroma_db"] = f"error: {str(e)[:50]}"
    
    # Check Ollama LLM
    try:
        if alfred_core and alfred_core.llm:
            components["ollama_llm"] = f"healthy (model: {alfred_core.model_name})"
        else:
            health_status = "degraded"
            components["ollama_llm"] = "not_available"
    except Exception as e:
        health_status = "degraded"
        components["ollama_llm"] = f"error: {str(e)[:50]}"
    
    # Check GPU
    if alfred_core:
        gpu_mgr = alfred_core.gpu_manager
        gpu_status_data = GPUStatus(
            gpu_available=gpu_mgr.has_gpu,
            device_type=gpu_mgr.device_type,
            device=gpu_mgr.device,
            gpu_info=gpu_mgr.gpu_info,
            memory_usage=gpu_mgr.get_memory_usage()
        )
        
        if gpu_mgr.has_gpu:
            components["gpu"] = f"available ({gpu_mgr.device_type})"
        else:
            components["gpu"] = "cpu_fallback"
    else:
        components["gpu"] = "unknown"
    
    return HealthResponse(
        status=health_status,
        timestamp=datetime.now().isoformat(),
        components=components,
        alfred_core_initialized=alfred_core is not None and alfred_core.is_initialized(),
        vectorstore_loaded=alfred_core.is_initialized() if alfred_core else False,
        gpu_status=gpu_status_data,
        is_fully_initialized=_initialization_complete
    )

@app.post("/query", response_model=QueryResponse, tags=["Consultas"])
async def query_alfred(request: QueryRequest):
    """
    Realizar una consulta a Alfred
    
    **NOTA DE SEGURIDAD**: Los datos personales se cifran automaticamente antes de almacenarlos
    y se descifran antes de enviarlos al cliente.
    
    **CIFRADO END-TO-END**: Las preguntas pueden venir cifradas desde el frontend y se descifran
    aqui antes de procesarlas.
    
    - **question**: Pregunta del usuario (puede estar cifrada con Fernet)
    - **use_history**: Buscar primero en el historial de respuestas
    - **save_response**: Guardar automáticamente la respuesta en el historial (con cifrado)
    - **search_documents**: Buscar en documentos o solo usar el prompt
    - **search_kwargs**: Parámetros adicionales para la búsqueda (k, fetch_k, etc.)
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        # DESCIFRAR PREGUNTA SI VIENE CIFRADA (End-to-End Encryption)
        question = request.question
        was_encrypted = False
        
        if question and question.startswith('gAAAAAB'):
            # La pregunta viene cifrada, descifrarla
            backend_logger.info(f"🔓 Pregunta cifrada detectada, descifrando...")
            try:
                question = decrypt_data(question)
                was_encrypted = True
                backend_logger.info(f"✅ Pregunta descifrada correctamente: {question[:50]}...")
            except Exception as decrypt_error:
                backend_logger.error(f"❌ Error al descifrar pregunta: {decrypt_error}")
                raise HTTPException(
                    status_code=400, 
                    detail="Error al descifrar la pregunta. Verifica que el cifrado este configurado correctamente."
                )
        
        print(f"Procesando consulta {'(descifrada)' if was_encrypted else ''}: {question[:50]}...")
        
        # Ejecutar consulta con version async optimizada
        result = await alfred_core.query_async(
            question=question,  # Usar pregunta descifrada
            use_history=request.use_history,
            search_documents=request.search_documents,
            search_kwargs=request.search_kwargs
        )
        
        print(f"Consulta procesada exitosamente")
        
        # Asegurar que los datos personales esten descifrados para el cliente
        personal_data = ensure_personal_data_decrypted(result.get('personal_data'))
        
        # Registrar acceso a datos personales si existen
        if personal_data:
            log_personal_data_access(
                operation="read",
                data_keys=list(personal_data.keys()),
                user_context=f"Query: {request.question[:30]}..."
            )
        
        # Guardar en historial si se solicita (se cifra automaticamente)
        if request.save_response and not result.get('from_history', False):
            functionsToHistory.save_qa_to_history(
                question=request.question,
                answer=result['answer'],
                personal_data=personal_data,
                sources=result.get('sources', []),
                encrypt_sensitive=True  # Cifrado explicito
            )
            
            if personal_data:
                log_personal_data_access(
                    operation="write",
                    data_keys=list(personal_data.keys()),
                    user_context="Guardado en historial"
                )
        
        # Crear respuesta
        response_data = QueryResponse(
            answer=result['answer'],
            personal_data=personal_data,
            sources=result.get('sources', []),
            from_history=result.get('from_history', False),
            history_score=result.get('history_score'),
            context_count=result.get('context_count', 0)
        )
        
        # CIFRAR DATOS SENSIBLES PARA VIAJE POR LA RED
        if is_encryption_enabled():
            response_dict = response_data.dict()
            response_dict = encrypt_for_transport(response_dict)
            response_data = QueryResponse(**response_dict)
        
        return response_data
    
    except Exception as e:
        # Sanitizar el mensaje de error para evitar problemas de encoding
        error_msg = str(e).encode('ascii', 'ignore').decode('ascii')
        if not error_msg:
            error_msg = "Error desconocido al procesar la consulta"
        
        print(f"Error al procesar consulta: {error_msg}")
        raise HTTPException(status_code=500, detail=f"Error al procesar consulta: {error_msg}")

@app.post("/history/search", response_model=List[HistoryEntry], tags=["Historial"])
async def search_history(request: HistorySearchRequest):
    """
    Buscar en el historial de preguntas y respuestas
    
    **NOTA DE SEGURIDAD**: Los datos personales se descifran automaticamente antes de ser enviados.
    
    - **search_term**: Término a buscar
    - **threshold**: Umbral de similitud (0.0 - 1.0)
    - **top_k**: Número máximo de resultados
    """
    try:
        # Buscar en historial (descifra automaticamente)
        results = functionsToHistory.search_in_qa_history(
            question=request.search_term,
            threshold=request.threshold,
            top_k=request.top_k
        )
        
        history_entries = []
        for score, entry in results:
            personal_data = ensure_personal_data_decrypted(entry.get('personal_data'))
            
            # Registrar acceso si hay datos personales
            if personal_data:
                log_personal_data_access(
                    operation="read",
                    data_keys=list(personal_data.keys()),
                    user_context=f"Busqueda en historial: {request.search_term[:30]}..."
                )
            
            history_entries.append(
                HistoryEntry(
                    timestamp=entry['timestamp'],
                    question=entry['question'],
                    answer=entry['answer'],
                    personal_data=personal_data,
                    sources=entry.get('sources', []),
                    similarity_score=score
                )
            )
        
        return history_entries
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al buscar en historial: {str(e)}")

@app.get("/history", response_model=List[HistoryEntry], tags=["Historial"])
async def get_history(limit: int = 10, offset: int = 0):
    """
    Obtener el historial de preguntas y respuestas
    
    **NOTA DE SEGURIDAD**: Los datos personales se descifran automaticamente antes de ser enviados.
    
    - **limit**: Número máximo de entradas a devolver
    - **offset**: Número de entradas a saltar (para paginación)
    """
    try:
        # Cargar historial (descifra automaticamente)
        history = functionsToHistory.load_qa_history(decrypt_sensitive=True)
        
        # Aplicar paginación
        paginated = history[offset:offset + limit]
        
        history_entries = []
        for entry in reversed(paginated):  # Más recientes primero
            personal_data = ensure_personal_data_decrypted(entry.get('personal_data'))
            
            # Registrar acceso si hay datos personales
            if personal_data:
                log_personal_data_access(
                    operation="read",
                    data_keys=list(personal_data.keys()),
                    user_context="Listado de historial"
                )
            
            history_entries.append(
                HistoryEntry(
                    timestamp=entry['timestamp'],
                    question=entry['question'],
                    answer=entry['answer'],
                    personal_data=personal_data,
                    sources=entry.get('sources', [])
                )
            )
        
        return history_entries
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener historial: {str(e)}")

@app.delete("/history/{timestamp}", tags=["Historial"])
async def delete_history_entry(timestamp: str):
    """
    Eliminar una entrada del historial por su timestamp
    
    **NOTA DE SEGURIDAD**: La eliminacion de datos sensibles se registra en el log de auditoria.
    
    - **timestamp**: Timestamp ISO de la entrada a eliminar
    """
    try:
        # Cargar entrada antes de eliminar para logging
        history = functionsToHistory.load_qa_history(decrypt_sensitive=False)
        entry_to_delete = next((e for e in history if e.get('timestamp') == timestamp), None)
        
        if entry_to_delete and entry_to_delete.get('personal_data'):
            log_personal_data_access(
                operation="delete",
                data_keys=list(entry_to_delete['personal_data'].keys()),
                user_context=f"Eliminacion de entrada: {timestamp}"
            )
        
        success = functionsToHistory.delete_qa_from_history(timestamp)
        
        if success:
            return {"status": "success", "message": "Entrada eliminada del historial"}
        else:
            raise HTTPException(status_code=404, detail="No se encontro la entrada con ese timestamp")
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al eliminar del historial: {str(e)}")

@app.get("/optimizations/stats", response_model=OptimizationStats, tags=["Optimizaciones"])
async def get_optimization_stats():
    """
    Obtener estadisticas de optimizaciones RAG
    
    Retorna informacion sobre:
    - Modelo de embeddings seleccionado automaticamente
    - Estadisticas de cache LRU (hits, misses, hit rate)
    - Estrategias de chunking aplicadas
    - Almacenamiento optimizado DuckDB+Parquet
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        # Obtener info del embedding manager
        embedding_manager = alfred_core._embedding_manager
        embedding_model = embedding_manager.select_best_model()
        vram_gb = embedding_manager.get_available_vram()
        
        # Dimension segun modelo
        model_dims = {
            "nomic-embed-text:v1.5": 768,
            "bge-large-en-v1.5": 1024,
            "gte-small": 384,
            "all-minilm:l6-v2": 384
        }
        embedding_dim = model_dims.get(embedding_model, 768)
        
        # Estadisticas de cache
        cache_stats = {"hits": 0, "misses": 0, "hit_rate": 0.0, "size": 0}
        if hasattr(alfred_core, '_retrieval_cache') and alfred_core._retrieval_cache:
            cache_stats = alfred_core._retrieval_cache.get_stats()
        
        # Info del vector manager
        vector_manager = alfred_core.vector_manager
        total_docs = 0
        if vector_manager and vector_manager._vectorstore:
            try:
                # Intentar obtener conteo de documentos
                total_docs = vector_manager._vectorstore._collection.count()
            except:
                total_docs = 0
        
        # Estrategias de chunking (placeholder - requiere tracking en chunking_manager)
        chunking_strategies = {
            "text": 0,
            "code": 0, 
            "document": 0
        }
        
        return OptimizationStats(
            embedding_model=embedding_model,
            embedding_dimension=embedding_dim,
            vram_available=vram_gb,
            cache_enabled=cache_stats.get("size", 0) >= 0,
            cache_hits=cache_stats.get("hits", 0),
            cache_misses=cache_stats.get("misses", 0),
            cache_hit_rate=cache_stats.get("hit_rate", 0.0),
            cache_size=cache_stats.get("size", 0),
            total_documents_indexed=total_docs,
            chunking_strategies=chunking_strategies,
            optimized_storage=vector_manager.use_optimized_storage if vector_manager else False,
            storage_path=vector_manager.chroma_db_path if vector_manager else ""
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener estadisticas de optimizaciones: {str(e)}")

@app.post("/reload", tags=["Mantenimiento"])
async def reload_documents(background_tasks: BackgroundTasks):
    """
    Recargar documentos desde el directorio configurado
    (Operación pesada, se ejecuta en background)
    """
    if not alfred_core:
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        # Ejecutar recarga en segundo plano
        background_tasks.add_task(alfred_core.reload_documents)
        return {
            "status": "started",
            "message": "Recarga de documentos iniciada en segundo plano"
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al recargar documentos: {str(e)}")

@app.get("/documents/test", tags=["Desarrollo"])
async def test_search(query: str, k: int = 5):
    """
    Realizar una búsqueda directa en la base de datos vectorial (para testing)
    
    - **query**: Consulta de búsqueda
    - **k**: Número de resultados a devolver
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        results = alfred_core.test_search(query, k)
        return {
            "query": query,
            "results_count": len(results),
            "results": [
                {
                    "source": doc.metadata.get('source', 'Desconocido'),
                    "content_preview": doc.page_content[:300] + "..." if len(doc.page_content) > 300 else doc.page_content
                }
                for doc in results
            ]
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al realizar búsqueda: {str(e)}")

@app.get("/model", response_model=ModelInfo, tags=["Configuración"])
async def get_current_model():
    """
    Obtener el modelo actual y los modelos disponibles
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        current_model = alfred_core.get_current_model()
        # Lista de modelos disponibles (puedes expandir esto según tus necesidades)
        available_models = ["gemma2:9b", "gpt-oss:20b", "gemma3:12b", "gemma3:4b", "gemma3:1b", "gemma3n:e4b"]
        
        return ModelInfo(
            model_name=current_model,
            available_models=available_models
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener información del modelo: {str(e)}")

@app.post("/model", tags=["Configuración"])
async def change_model(request: ChangeModelRequest):
    """
    Cambiar el modelo LLM actual (lazy loading)
    
    El modelo se configura pero NO se carga hasta la proxima consulta.
    Esto permite cambiar entre modelos instantaneamente sin esperas.
    
    - **model_name**: Nombre del nuevo modelo a utilizar
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        print(f"[BACKEND] Configurando modelo: {request.model_name}")
        success = alfred_core.change_model(request.model_name)
        
        if success:
            print(f"[BACKEND] Modelo configurado: {request.model_name}")
            return {
                "status": "success",
                "message": f"Modelo configurado como {request.model_name}. Se cargará en la próxima consulta.",
                "model_name": request.model_name
            }
        else:
            error_msg = f"No se pudo configurar el modelo {request.model_name}"
            print(f"[BACKEND] {error_msg}")
            raise HTTPException(status_code=500, detail=error_msg)
    
    except HTTPException:
        raise
    except Exception as e:
        error_detail = f"Error al configurar modelo: {str(e)}"
        print(f"[BACKEND ERROR] {error_detail}")
        raise HTTPException(status_code=500, detail=error_detail)

@app.get("/model/current", tags=["Configuración"])
async def get_current_model():
    """
    Obtener el modelo LLM actualmente configurado
    
    Retorna el modelo activo almacenado en base de datos SQLite.
    Este es el modelo que se cargara en la proxima consulta.
    
    Returns:
        {
            "model_name": "gemma2:9b",
            "source": "database" | "env" | "default",
            "is_loaded": true | false
        }
    """
    if not alfred_core:
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        from db_manager import get_model_setting
        
        # Modelo actualmente configurado en memoria
        current_model = alfred_core.get_current_model()
        
        # Modelo guardado en BD (puede ser diferente si acaba de iniciar)
        db_model = get_model_setting('last_used_model')
        
        # Verificar si LLM esta cargado en memoria
        is_loaded = alfred_core._llm is not None
        
        # Determinar fuente del modelo
        if db_model and db_model == current_model:
            source = "database"
        elif current_model == os.getenv('ALFRED_MODEL', 'gemma2:9b'):
            source = "env"
        else:
            source = "memory"
        
        return {
            "model_name": current_model,
            "source": source,
            "is_loaded": is_loaded,
            "db_model": db_model,
            "env_model": os.getenv('ALFRED_MODEL', 'gemma2:9b')
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener modelo actual: {str(e)}")

# ====================================
# ENDPOINTS DE MODO DE APLICACION
# ====================================

@app.post("/settings/mode", response_model=ModeResponse, tags=["Configuración"])
async def set_mode(request: ModeRequest):
    """
    Guardar el modo actual de la aplicacion
    
    Modos disponibles:
    - work: Modo trabajo (cyan)
    - focus: Modo concentracion (morado)
    - personal: Modo personal (rosa)
    - creative: Modo creativo (naranja)
    
    El modo se guarda en la base de datos y se carga automaticamente
    al iniciar la aplicacion.
    """
    try:
        from db_manager import set_user_setting
        from datetime import datetime
        
        # Guardar en base de datos
        success = set_user_setting('app_mode', request.mode, 'string')
        
        if not success:
            raise Exception("No se pudo guardar el modo en la base de datos")
        
        backend_logger.info(f"Modo cambiado a: {request.mode}")
        
        return ModeResponse(
            mode=request.mode,
            updated_at=datetime.now().isoformat()
        )
        
    except Exception as e:
        error_detail = f"Error al guardar modo: {str(e)}"
        backend_logger.error(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)

@app.get("/settings/mode", response_model=ModeResponse, tags=["Configuración"])
async def get_mode():
    """
    Obtener el modo actual de la aplicacion
    
    Retorna el modo guardado en la base de datos.
    Si no hay modo guardado, retorna 'work' como valor por defecto.
    """
    try:
        from db_manager import get_user_setting
        
        # Obtener de base de datos (retorna directamente el valor string)
        mode = get_user_setting('app_mode', default='work')
        
        return ModeResponse(
            mode=mode,
            updated_at=None  # Opcional: podriamos consultar updated_at si es necesario
        )
        
    except Exception as e:
        error_detail = f"Error al obtener modo: {str(e)}"
        backend_logger.error(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)

# ====================================
# ENDPOINTS DE TEMA DE APLICACION
# ====================================

@app.post("/settings/theme", response_model=ThemeResponse, tags=["Configuración"])
async def set_theme(request: ThemeRequest):
    """
    Guardar el tema actual de la aplicacion
    
    Temas disponibles:
    - light: Tema claro
    - dark: Tema oscuro
    
    El tema se guarda en la base de datos y se carga automaticamente
    al iniciar la aplicacion.
    """
    try:
        from db_manager import set_user_setting
        from datetime import datetime
        
        # Guardar en base de datos
        success = set_user_setting('app_theme', request.theme, 'string')
        
        if not success:
            raise Exception("No se pudo guardar el tema en la base de datos")
        
        backend_logger.info(f"Tema cambiado a: {request.theme}")
        
        return ThemeResponse(
            theme=request.theme,
            updated_at=datetime.now().isoformat()
        )
        
    except Exception as e:
        error_detail = f"Error al guardar tema: {str(e)}"
        backend_logger.error(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)

@app.get("/settings/theme", response_model=ThemeResponse, tags=["Configuración"])
async def get_theme():
    """
    Obtener el tema actual de la aplicacion
    
    Retorna el tema guardado en la base de datos.
    Si no hay tema guardado, retorna 'dark' como valor por defecto.
    """
    try:
        from db_manager import get_user_setting
        
        # Obtener de base de datos (retorna directamente el valor string)
        theme = get_user_setting('app_theme', default='dark')
        
        return ThemeResponse(
            theme=theme,
            updated_at=None  # Opcional: podriamos consultar updated_at si es necesario
        )
        
    except Exception as e:
        error_detail = f"Error al obtener tema: {str(e)}"
        backend_logger.error(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)

@app.get("/gpu/status", response_model=GPUStatus, tags=["Sistema"])
async def get_gpu_status():
    """
    Obtener el estado actual de la GPU
    
    Retorna información detallada sobre:
    - Si hay GPU disponible
    - Tipo de GPU (NVIDIA, AMD, Apple Silicon)
    - Uso de memoria (si está disponible)
    - Información del dispositivo
    """
    if not alfred_core:
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        gpu_mgr = alfred_core.gpu_manager
        return GPUStatus(
            gpu_available=gpu_mgr.has_gpu,
            device_type=gpu_mgr.device_type,
            device=gpu_mgr.device,
            gpu_info=gpu_mgr.gpu_info,
            memory_usage=gpu_mgr.get_memory_usage()
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener estado de GPU: {str(e)}")

@app.get("/gpu/report", tags=["Sistema"])
async def get_gpu_report():
    """
    Obtener un reporte detallado del estado de la GPU en formato de texto
    """
    if not alfred_core:
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    try:
        gpu_mgr = alfred_core.gpu_manager
        report = gpu_mgr.get_status_report()
        return {
            "report": report,
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al generar reporte: {str(e)}")

@app.post("/ollama/stop", tags=["Configuración"])
async def stop_ollama(background_tasks: BackgroundTasks):
    """
    Detener el servicio de Ollama manualmente para liberar recursos
    
    Ejecuta el comando 'ollama stop <modelo>' para descargar el modelo de la memoria
    Se ejecuta en segundo plano para no bloquear la respuesta
    """
    import subprocess
    
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no está inicializado")
    
    # Obtener el modelo actual
    current_model = alfred_core.get_current_model()
    
    # Función para ejecutar en segundo plano
    def stop_ollama_background():
        try:
            print(f"[Background] Intentando detener modelo: {current_model}")
            
            # Ejecutar comando ollama stop con el nombre del modelo
            result = subprocess.run(
                ["ollama", "stop", current_model],
                capture_output=True,
                text=True,
                timeout=10
            )
            
            print(f"[Background] Resultado del comando: returncode={result.returncode}")
            print(f"[Background] stdout: {result.stdout}")
            print(f"[Background] stderr: {result.stderr}")
            
            if result.returncode == 0:
                print(f"[Background] Modelo {current_model} detenido exitosamente")
            else:
                print(f"[Background] Error al detener modelo: {result.stderr}")
                
        except Exception as e:
            print(f"[Background] Excepción al detener Ollama: {str(e)}")
    
    # Agregar tarea en segundo plano
    background_tasks.add_task(stop_ollama_background)
    
    # Responder inmediatamente sin esperar
    return {
        "status": "success",
        "message": f"Deteniendo modelo {current_model} en segundo plano. Los recursos se liberaran en unos segundos.",
        "model": current_model
    }

@app.get("/ollama/keep-alive", response_model=OllamaKeepAliveResponse, tags=["Configuración"])
async def get_ollama_keep_alive():
    """
    Obtener el tiempo de keep_alive actual de Ollama
    
    El keep_alive controla cuanto tiempo Ollama mantiene el modelo en memoria
    despues de usarlo. Despues de este tiempo sin uso, el modelo se descarga
    automaticamente para liberar recursos.
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no esta inicializado")
    
    try:
        keep_alive = alfred_core.get_ollama_keep_alive()
        
        if keep_alive == 0:
            description = "El modelo se descarga inmediatamente despues de cada uso"
        elif keep_alive < 60:
            description = f"El modelo permanece en memoria {keep_alive} segundos despues de cada uso"
        elif keep_alive < 3600:
            minutes = keep_alive // 60
            description = f"El modelo permanece en memoria {minutes} minuto(s) despues de cada uso"
        else:
            hours = keep_alive // 3600
            description = f"El modelo permanece en memoria {hours} hora(s) despues de cada uso"
        
        return OllamaKeepAliveResponse(
            keep_alive_seconds=keep_alive,
            description=description
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al obtener keep_alive: {str(e)}")

@app.put("/ollama/keep-alive", tags=["Configuración"])
async def set_ollama_keep_alive(request: OllamaKeepAliveRequest):
    """
    Actualizar el tiempo de keep_alive de Ollama
    
    - **seconds**: Tiempo en segundos (0-3600)
        - 0: Descargar el modelo inmediatamente despues de cada uso
        - 1-3600: Mantener el modelo en memoria durante este tiempo
    
    Valores recomendados:
    - 0s: Para liberar memoria inmediatamente (uso ocasional)
    - 30s-60s: Para uso regular con pausas cortas
    - 300s (5min): Para sesiones de trabajo activas
    - 1800s (30min): Para uso prolongado sin pausas
    """
    if not alfred_core or not alfred_core.is_initialized():
        raise HTTPException(status_code=503, detail="Alfred Core no esta inicializado")
    
    try:
        success = alfred_core.set_ollama_keep_alive(request.seconds)
        
        if not success:
            raise HTTPException(status_code=400, detail="No se pudo actualizar el keep_alive")
        
        return {
            "status": "success",
            "message": f"keep_alive actualizado a {request.seconds} segundos",
            "keep_alive_seconds": request.seconds
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al actualizar keep_alive: {str(e)}")

# ====================================
# ENDPOINTS DE GESTION DE MODELOS OLLAMA
# ====================================

@app.get("/ollama/models/list", tags=["Modelos"])
async def list_ollama_models():
    """
    Listar todos los modelos instalados en Ollama
    
    Ejecuta 'ollama list' y retorna la lista de modelos disponibles
    con sus detalles (nombre, tamaño, fecha de modificacion)
    """
    import subprocess
    import json
    
    try:
        # Ejecutar comando ollama list
        result = subprocess.run(
            ["ollama", "list"],
            capture_output=True,
            text=True,
            timeout=10
        )
        
        if result.returncode != 0:
            raise HTTPException(
                status_code=500, 
                detail=f"Error al ejecutar ollama list: {result.stderr}"
            )
        
        # Parsear la salida
        lines = result.stdout.strip().split('\n')
        
        if len(lines) < 1:  # Al menos un modelo
            return {
                "models": [],
                "count": 0,
                "message": "No hay modelos instalados"
            }
        
        # Procesar todas las lineas (saltar header si existe)
        models = []
        for line in lines:
            # Limpiar la linea y dividir por multiples espacios
            line = line.strip()
            if not line:
                continue
            
            # Separar por espacios multiples (usa split sin argumentos)
            parts = line.split()
            
            # Saltar el header (primera linea con NAME, ID, SIZE, MODIFIED)
            if parts[0] == 'NAME' or parts[0] == 'name':
                continue
            
            if len(parts) >= 4:
                model_name = parts[0]
                model_id = parts[1]
                # El tamano son 2 partes: numero + unidad (ej: 669 MB)
                size = f"{parts[2]} {parts[3]}" if len(parts) > 3 else parts[2]
                # Los ultimos elementos son la fecha/tiempo (despues del tamano)
                modified = ' '.join(parts[4:]) if len(parts) > 4 else ""
                
                models.append({
                    "name": model_name,
                    "id": model_id,
                    "size": size,
                    "modified": modified
                })
        
        backend_logger.info(f"Modelos Ollama encontrados: {len(models)}")
        
        return {
            "models": models,
            "count": len(models),
            "timestamp": datetime.now().isoformat()
        }
        
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=504, detail="Timeout al listar modelos de Ollama")
    except Exception as e:
        backend_logger.error(f"Error al listar modelos: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al listar modelos: {str(e)}")

@app.post("/ollama/models/download", tags=["Modelos"])
async def download_ollama_model(background_tasks: BackgroundTasks, model_name: str):
    """
    Descargar un modelo de Ollama
    
    Ejecuta 'ollama pull <model_name>' en segundo plano.
    El proceso puede tardar varios minutos dependiendo del tamaño del modelo.
    
    - **model_name**: Nombre del modelo a descargar (ej: llama2, mistral, gemma2:9b)
    
    La descarga se ejecuta en segundo plano y el endpoint retorna inmediatamente.
    """
    import subprocess
    import re
    from db_manager import save_model_download_history, update_model_download_progress
    
    if not model_name or model_name.strip() == "":
        raise HTTPException(status_code=400, detail="El nombre del modelo es requerido")
    
    # Función para ejecutar en segundo plano
    def download_model_background():
        import traceback
        try:
            print(f"[DEBUG] Iniciando descarga del modelo: {model_name}")
            backend_logger.info(f"Iniciando descarga del modelo: {model_name}")
            
            # Guardar inicio de descarga en BD
            try:
                print(f"[DEBUG] Guardando en BD: {model_name}")
                save_model_download_history(model_name, 'downloading', 'Descarga iniciada', 0)
                print(f"[DEBUG] Guardado exitoso en BD")
            except Exception as e:
                print(f"[DEBUG ERROR] Error al guardar historial: {str(e)}")
                print(f"[DEBUG ERROR] Traceback: {traceback.format_exc()}")
                backend_logger.error(f"Error al guardar historial: {str(e)}")
            except Exception as e:
                backend_logger.error(f"Error al guardar historial: {str(e)}")
            
            # Ejecutar comando ollama pull con captura de output en tiempo real
            process = subprocess.Popen(
                ["ollama", "pull", model_name],
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                bufsize=1,
                universal_newlines=True,
                encoding='utf-8',
                errors='replace'  # Manejar caracteres especiales de Windows
            )
            
            # Leer output linea por linea para capturar progreso
            last_progress = 0
            last_clean_message = ""
            
            for line in process.stdout:
                # Limpieza simple de caracteres de control
                # Remover códigos ANSI comunes: [K, [A, [1G, [?25h, etc
                clean_line = line
                for code in ['[K', '[A', '[1G', '[?25h', '[?25l', '[?2026l', '[?2026h']:
                    clean_line = clean_line.replace(code, '')
                
                # Remover otros caracteres de control
                clean_line = ''.join(char for char in clean_line if char.isprintable() or char in '\n\r\t')
                clean_line = clean_line.strip()
                
                if not clean_line or clean_line == 'pulling manifest':
                    continue
                
                # Intentar extraer porcentaje de progreso
                if '%' in clean_line:
                    try:
                        # Buscar el porcentaje (patrón simple)
                        percent_pos = clean_line.find('%')
                        if percent_pos > 0:
                            # Buscar hacia atrás hasta encontrar el número
                            start = percent_pos - 1
                            while start >= 0 and (clean_line[start].isdigit() or clean_line[start] == ' '):
                                start -= 1
                            progress_str = clean_line[start+1:percent_pos].strip()
                            
                            if progress_str.isdigit():
                                progress = int(progress_str)
                                
                                # Crear mensaje simple sin caracteres especiales
                                # Remover bloques visuales █ ▕ ▏
                                simple_message = clean_line.replace('█', '').replace('▕', '').replace('▏', '')
                                simple_message = ' '.join(simple_message.split())  # Normalizar espacios
                                
                                # Limitar longitud del mensaje
                                if len(simple_message) > 100:
                                    simple_message = simple_message[:97] + '...'
                                
                                if progress != last_progress:
                                    last_progress = progress
                                    last_clean_message = simple_message
                                    update_model_download_progress(model_name, progress, simple_message)
                                    print(f"[DEBUG] Progreso: {progress}%")
                                    backend_logger.info(f"Progreso: {progress}%")
                    except Exception as e:
                        print(f"[DEBUG] Error parseando: {str(e)}")
                        pass
            
            # Esperar a que termine el proceso
            return_code = process.wait()
            
            print(f"[DEBUG] Proceso terminado con código: {return_code}")
            
            if return_code == 0:
                print(f"[DEBUG] Descarga exitosa, guardando en BD como completed...")
                backend_logger.info(f"Modelo {model_name} descargado exitosamente")
                
                # Actualizar BD con exito
                try:
                    save_model_download_history(model_name, 'completed', 'Descarga completada', 100)
                    print(f"[DEBUG] Estado 'completed' guardado exitosamente")
                except Exception as e:
                    print(f"[DEBUG ERROR] Error al guardar completed: {str(e)}")
                    backend_logger.error(f"Error al actualizar historial: {str(e)}")
                    import traceback
                    print(f"[DEBUG ERROR] Traceback: {traceback.format_exc()}")
            else:
                error_msg = f"Error al descargar modelo (return code: {return_code})"
                print(f"[DEBUG] Error en descarga: {error_msg}")
                backend_logger.error(f"Error al descargar modelo {model_name}: {error_msg}")
                
                # Guardar error en BD
                try:
                    save_model_download_history(model_name, 'failed', error_msg, last_progress)
                except Exception as e:
                    backend_logger.error(f"Error al guardar error en historial: {str(e)}")
                
        except Exception as e:
            backend_logger.error(f"Excepcion al descargar modelo {model_name}: {str(e)}")
            try:
                save_model_download_history(model_name, 'failed', str(e), 0)
            except:
                pass
    
    # Agregar tarea en segundo plano
    background_tasks.add_task(download_model_background)
    
    # Responder inmediatamente
    return {
        "status": "downloading",
        "message": f"Descarga de {model_name} iniciada en segundo plano. Esto puede tardar varios minutos.",
        "model_name": model_name,
        "timestamp": datetime.now().isoformat()
    }

@app.delete("/ollama/models/{model_name}", tags=["Modelos"])
async def delete_ollama_model(model_name: str):
    """
    Eliminar un modelo de Ollama
    
    Ejecuta 'ollama rm <model_name>' para eliminar el modelo del sistema.
    
    - **model_name**: Nombre del modelo a eliminar (debe estar instalado)
    """
    import subprocess
    from db_manager import save_model_download_history
    
    if not model_name or model_name.strip() == "":
        raise HTTPException(status_code=400, detail="El nombre del modelo es requerido")
    
    try:
        backend_logger.info(f"Eliminando modelo: {model_name}")
        
        # Ejecutar comando ollama rm
        result = subprocess.run(
            ["ollama", "rm", model_name],
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode == 0:
            backend_logger.info(f"Modelo {model_name} eliminado exitosamente")
            
            # Guardar en historial
            try:
                save_model_download_history(model_name, 'deleted', 'Modelo eliminado')
            except Exception as e:
                backend_logger.error(f"Error al guardar en historial: {str(e)}")
            
            return {
                "status": "success",
                "message": f"Modelo {model_name} eliminado exitosamente",
                "model_name": model_name
            }
        else:
            error_msg = result.stderr or "Error desconocido"
            backend_logger.error(f"Error al eliminar modelo: {error_msg}")
            raise HTTPException(status_code=500, detail=error_msg)
            
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=504, detail="Timeout al eliminar modelo")
    except HTTPException:
        raise
    except Exception as e:
        backend_logger.error(f"Error al eliminar modelo: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al eliminar modelo: {str(e)}")

@app.get("/ollama/models/download-history", tags=["Modelos"])
async def get_model_download_history():
    """
    Obtener el historial de descargas de modelos
    
    Retorna todas las descargas registradas en la base de datos
    con su estado (downloading, completed, failed, deleted)
    """
    try:
        from db_manager import get_model_download_history
        
        history = get_model_download_history()
        
        return {
            "history": history,
            "count": len(history),
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        backend_logger.error(f"Error al obtener historial: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al obtener historial: {str(e)}")

@app.get("/ollama/models/status/{model_name}", tags=["Modelos"])
async def get_model_download_status_endpoint(model_name: str):
    """
    Obtener el estado actual de descarga de un modelo especifico
    
    Util para polling y obtener progreso en tiempo real durante la descarga.
    
    - **model_name**: Nombre del modelo
    
    Retorna el ultimo estado del modelo con su progreso (0-100)
    """
    try:
        from db_manager import get_model_download_status
        
        status = get_model_download_status(model_name)
        
        if status:
            return {
                "found": True,
                "model_name": model_name,
                "status": status['status'],
                "progress": status['progress'],
                "message": status['message'],
                "updated_at": status['updated_at']
            }
        else:
            return {
                "found": False,
                "model_name": model_name,
                "status": "not_found",
                "progress": 0,
                "message": "No hay informacion de descarga para este modelo"
            }
        
    except Exception as e:
        backend_logger.error(f"Error al obtener estado de modelo: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al obtener estado: {str(e)}")

# --- Manejo de errores global ---
@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    """
    Manejador global de excepciones mejorado
    
    Proporciona:
    - Error IDs unicos para tracking
    - Logging estructurado completo
    - Respuestas sanitizadas (sin stack traces sensibles)
    - Timestamp ISO para auditoria
    """
    import traceback
    
    # Generar ID unico para el error
    error_id = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    
    # Log completo del error con stack trace
    backend_logger.error(
        f"[{error_id}] Excepcion no manejada capturada",
        extra={
            "error_id": error_id,
            "error_type": type(exc).__name__,
            "error_message": str(exc),
            "request_path": str(request.url) if hasattr(request, 'url') else "unknown",
            "request_method": request.method if hasattr(request, 'method') else "unknown"
        }
    )
    backend_logger.error(f"[{error_id}] Stack trace:\n{traceback.format_exc()}")
    
    # Sanitizar mensaje de error para el cliente
    # Eliminar caracteres no-ASCII que puedan causar problemas
    error_msg = str(exc).encode('ascii', 'ignore').decode('ascii')
    if not error_msg or len(error_msg.strip()) == 0:
        error_msg = "Error interno del servidor"
    
    # Respuesta estructurada
    from fastapi.responses import JSONResponse
    return JSONResponse(
        status_code=500,
        content={
            "error": "Internal Server Error",
            "detail": error_msg,
            "error_id": error_id,
            "timestamp": datetime.now().isoformat(),
            "message": f"Error registrado con ID: {error_id}. Revisa los logs para mas detalles."
        }
    )

if __name__ == "__main__":
    import uvicorn
    
    # Obtener configuración desde variables de entorno
    host = os.getenv("ALFRED_HOST", "127.0.0.1")
    port = int(os.getenv("ALFRED_PORT", "8000"))
    # Auto-reload deshabilitado - usar modo desarrollo de uvicorn manualmente si es necesario
    reload = False
    
    print(f"""
    ============================================================
              Alfred Backend API Server
    ============================================================
      URL:    http://{host}:{port}
      Docs:   http://{host}:{port}/docs
      Health: http://{host}:{port}/health
    ============================================================
    """)
    
    uvicorn.run(
        "alfred_backend:app",
        host=host,
        port=port,
        reload=reload,
        log_level="info"
    )
